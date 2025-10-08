#!/usr/bin/env python3
"""
Document Processor for Study Space Application
Extracts content from PDF and PPT files, processes with Qwen3-coder LLM, and stores in database
"""

import os
import sys
import json
import logging
import argparse
from pathlib import Path
from typing import Dict, List, Optional, Tuple
import requests
import mysql.connector
from mysql.connector import Error
import PyPDF2
from pptx import Presentation
from docx import Document
import uuid
from datetime import datetime

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Main class for document processing and LLM integration"""
    
    def __init__(self, config_file: str = "../../config/ai-processor/config.json"):
        """Initialize the document processor with configuration"""
        self.config = self.load_all_configs()
        self.db_connection = None
        self.llm_endpoint = self.config.get('llm_endpoint', 'http://localhost:11434')
        self.llm_model = self.config.get('llm_model', 'qwen3-coder:480b-cloud')
        
    def load_all_configs(self) -> Dict:
        """Load all configuration files and merge them"""
        config_files = {
            'ai-processor': '../../config/ai-processor/config.json',
            'database': '../../config/database/config.json',
            'api': '../../config/api/config.json'
        }
        
        merged_config = {}
        
        for config_type, config_path in config_files.items():
            try:
                with open(config_path, 'r') as f:
                    config_data = json.load(f)
                    merged_config.update(config_data)
                    logger.info(f"Loaded {config_type} configuration from {config_path}")
            except FileNotFoundError:
                logger.warning(f"Config file {config_path} not found, using defaults for {config_type}")
                if config_type == 'database':
                    merged_config['database'] = {
                        "host": "localhost",
                        "port": 3306,
                        "user": "root",
                        "password": "",
                        "database": "studyspace_db"
                    }
                elif config_type == 'ai-processor':
                    merged_config.update({
                        'llm_endpoint': 'http://localhost:11434',
                        'llm_model': 'qwen3-coder:480b-cloud',
                        'output_dir': './processed_documents',
                        'supported_formats': ['.pdf', '.ppt', '.pptx', '.doc', '.docx']
                    })
                elif config_type == 'api':
                    merged_config['api'] = {
                        "host": "localhost",
                        "port": 8000,
                        "debug": True
                    }
            except Exception as e:
                logger.error(f"Error loading {config_type} configuration: {e}")
        
        return merged_config
    
    def load_config(self, config_file: str) -> Dict:
        """Load configuration from JSON file (legacy method)"""
        try:
            with open(config_file, 'r') as f:
                return json.load(f)
        except FileNotFoundError:
            logger.warning(f"Config file {config_file} not found, using defaults")
            return self.get_default_config()
    
    def get_default_config(self) -> Dict:
        """Return default configuration"""
        return {
            "database": {
                "host": "localhost",
                "port": 3306,
                "user": "root",
                "password": "",
                "database": "studyspace_db"
            },
            "llm_endpoint": "http://localhost:11434",
            "llm_model": "qwen3-coder:480b-cloud",
            "output_dir": "./processed_documents",
            "supported_formats": [".pdf", ".ppt", ".pptx", ".doc", ".docx"]
        }
    
    def connect_database(self) -> bool:
        """Connect to MySQL database"""
        try:
            db_config = self.config['database']
            self.db_connection = mysql.connector.connect(
                host=db_config['host'],
                port=db_config['port'],
                user=db_config['user'],
                password=db_config['password'],
                database=db_config['database']
            )
            logger.info("Successfully connected to database")
            return True
        except Error as e:
            logger.error(f"Database connection failed: {e}")
            return False
    
    def extract_pdf_content(self, file_path: str) -> str:
        """Extract text content from PDF file"""
        try:
            content = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    content += page.extract_text() + "\n"
            logger.info(f"Extracted {len(content)} characters from PDF: {file_path}")
            return content.strip()
        except Exception as e:
            logger.error(f"Error extracting PDF content: {e}")
            return ""
    
    def extract_ppt_content(self, file_path: str) -> str:
        """Extract text content from PowerPoint file"""
        try:
            content = ""
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides):
                content += f"\n--- Slide {slide_num + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
            
            logger.info(f"Extracted {len(content)} characters from PPT: {file_path}")
            return content.strip()
        except Exception as e:
            logger.error(f"Error extracting PPT content: {e}")
            return ""
    
    def extract_docx_content(self, file_path: str) -> str:
        """Extract text content from Word document"""
        try:
            doc = Document(file_path)
            content = ""
            for paragraph in doc.paragraphs:
                content += paragraph.text + "\n"
            
            logger.info(f"Extracted {len(content)} characters from DOCX: {file_path}")
            return content.strip()
        except Exception as e:
            logger.error(f"Error extracting DOCX content: {e}")
            return ""
    
    def extract_content(self, file_path: str) -> Tuple[str, str]:
        """Extract content from supported file types"""
        file_path = Path(file_path)
        file_extension = file_path.suffix.lower()
        
        if file_extension == '.pdf':
            content = self.extract_pdf_content(str(file_path))
            file_type = "PDF"
        elif file_extension in ['.ppt', '.pptx']:
            content = self.extract_ppt_content(str(file_path))
            file_type = "PowerPoint"
        elif file_extension in ['.doc', '.docx']:
            content = self.extract_docx_content(str(file_path))
            file_type = "Word Document"
        else:
            logger.error(f"Unsupported file type: {file_extension}")
            return "", ""
        
        return content, file_type
    
    def process_with_llm(self, content: str, file_type: str) -> Dict:
        """Process extracted content with Qwen3-coder LLM"""
        try:
            # Prepare the prompt for the LLM
            prompt = f"""
            You are an intelligent educational content processor. Analyze the imported text carefully and create both comprehensive study notes and educational flashcards.

            Content to analyze:
            {content[:80000]}  # Maximum content limit for complete document analysis
            
            For STUDY NOTES, act as an educational content analyzer:
            You are an academic document analyzer designed to create detailed educational outlines and study guides from any imported text file.

            Read the entire document carefully, then generate a structured, comprehensive outline summarizing the most important topics, definitions, theories, causes/effects, and related laws or examples.

            Your goals:

            Identify and clearly present all main topics and subtopics.

            Summarize key concepts using short, clear bullet points or concise paragraphs.

            Extract definitions, theories, historical developments, and laws/regulations relevant to the document.

            Maintain a logical and academic flow (overview → development → key concepts → laws → issues → causes/effects → conclusion).

            Use clear section headers and a readable study-guide style format.

            Output Format Example
            Title: [Auto-detect or summarize the title]

            I. Overview  
            - Brief summary of the document's main topic or theme.

            II. Historical Development  
            - Key milestones, events, or changes related to the topic.  

            III. Key Terms / Concepts  
            - [Term]: [Concise definition or explanation]  
            - Environmental Awareness: A sense of responsibility towards preserving the environment for a sustainable future.  
            - Environmentalism: A movement aimed at protecting the environment from human impacts.  
            - Pro-Environmental Behavior (PEB): Actions individuals take to protect the environment.  
            - Green Consumerism: Buying eco-friendly and sustainable products.

            IV. Fundamental Theories / Models  
            | Theory/Model | Description |  
            |---------------|-------------|  
            | Ecological Awareness | Understanding the impact of human actions on nature. |  
            | Sustainable Development | Meeting present needs without harming future generations. |  
            | Climate Change | Long-term shifts in weather patterns caused by human activities. |  

            V. Key Historical Developments  
            - Indus Civilization (3000 BC): Established waste management standards.  
            - Industrial Revolution: Beginning of formal environmentalism due to pollution.  
            - Clean Air Act (1956): Reduced air pollution and improved public health.  

            VI. Key Regulations / Legislation  
            | Regulation/Legislation | Description |  
            |-------------------------|-------------|  
            | Presidential Decree 1586 | Establishes the Environmental Impact Statement System. |  
            | RA 6969 | Toxic Substances and Hazardous Waste Act (1990). |  
            | RA 849 | Philippine Clean Air Act (1999). |  
            | RA 9003 | Ecological Solid Waste Management Act (2000). |  
            | RA 9275 | Clean Water Act (2004). |  

            VII. Key Environmental Issues  
            - Climate Change: Long-term global temperature rise.  
            - Stratospheric Ozone Depletion: Thinning of ozone layer from chemicals like CFCs.  
            - Acid Deposition: Acid rain harming soil and ecosystems.  

            VIII. Facts to Memorize  
            - Environmentalism aims to reduce human environmental impact.  
            - The ozone layer protects Earth from harmful UV radiation.  
            - Thermal inversion causes smog and air pollution near the ground.  

            IX. Cause and Effect Relationships  
            | Cause | Effect |  
            |--------|--------|  
            | Industrial Revolution | Pollution → Birth of environmentalism. |  
            | Deforestation | Climate change and biodiversity loss. |  
            | Use of CFCs | Ozone layer depletion. |  
            | Burning fossil fuels | Acid rain formation and soil damage. |  
            | Thermal inversion | Air pollution trapped → health problems. |  

            X. Reference Information  
            - Clean Air Act: Regulates air pollution globally.  
            - WWF and RSPB: Organizations focused on environmental protection.  
            - Green Consumerism: Choosing sustainable brands and products.  

            XI. Conclusion / Future Directions  
            - Summary of what can be done for sustainability.  
            - Role of education, activism, and technology in protecting the environment.  


            Output Requirements:

            Format sections in a clean, academic style.

            Use bullet points or tables for clarity.

            Ensure no information is lost from the source text, only summarized.

            Maintain educational tone (like a study guide or review sheet).
            
            For FLASHCARDS, act as an intelligent flashcard generator:
            - Identify key points – Extract important definitions, dates, laws, processes, causes and effects, examples, and important people or events
            - Simplify and clarify – Rephrase complex ideas into short, clear, and direct questions and answers
            - Promote understanding – Focus not only on memorization but also on conceptual comprehension (e.g., "Why" and "How" questions)
            - Avoid redundancy – Each flashcard should cover one key concept only
            
            Format your response as JSON with the following structure:
            {{
                "title": "string",
                "summary": "string",
                "key_topics": ["topic1", "topic2", ...],
                "definitions": ["term1: definition1", "term2: definition2", ...],
                "study_notes": "Title: [Auto-detect or summarize the title]\\n\\nI. Overview\\n- Brief summary of the document's main topic or theme.\\n\\nII. Historical Development\\n- Key milestones, events, or changes related to the topic.\\n\\nIII. Key Terms / Concepts\\n- [Term]: [Concise definition or explanation]\\n\\nIV. Fundamental Theories / Models\\n| Theory/Model | Description |\\n|---------------|-------------|\\n\\nV. Key Historical Developments\\n- [Event/Date]: [Description]\\n\\nVI. Key Regulations / Legislation\\n| Regulation/Legislation | Description |\\n|-------------------------|-------------|\\n\\nVII. Key Environmental Issues\\n- [Issue]: [Description]\\n\\nVIII. Facts to Memorize\\n- [Important fact]\\n\\nIX. Cause and Effect Relationships\\n| Cause | Effect |\\n|--------|--------|\\n\\nX. Reference Information\\n- [Reference]: [Description]\\n\\nXI. Conclusion / Future Directions\\n- Summary of what can be done for sustainability.\\n- Role of education, activism, and technology in protecting the environment.",
                "difficulty": "Easy|Medium|Hard",
                "subject": "string",
                "flashcards": [
                    {{
                        "question": "What is [concept]?",
                        "answer": "Clear, concise explanation (1-3 sentences)"
                    }},
                    {{
                        "question": "When did [event] occur?",
                        "answer": "Specific date and context"
                    }},
                    {{
                        "question": "What does [law/principle] state?",
                        "answer": "Brief statement of the law with key points"
                    }},
                    {{
                        "question": "What causes [phenomenon]?",
                        "answer": "List of main causes with brief explanations"
                    }},
                    {{
                        "question": "Give an example of [concept]",
                        "answer": "Specific real-world example with brief explanation"
                    }}
                ]
            }}
            
            Generate 20-50 flashcards covering different aspects of the content. Focus on creating meaningful, educational questions that promote understanding rather than simple memorization. Ensure comprehensive coverage of all topics in the document.
            """
            
            # Call the LLM API
            response = requests.post(
                f"{self.llm_endpoint}/api/generate",
                json={
                    "model": self.llm_model,
                    "prompt": prompt,
                    "stream": False,
                    "options": {
                        "temperature": 0.3,
                        "top_p": 0.9,
                        "max_tokens": 40000
                    }
                },
                timeout=300
            )
            
            if response.status_code == 200:
                result = response.json()
                llm_response = result.get('response', '')
                
                # Try to parse JSON response
                try:
                    # Extract JSON from the response (LLM might add extra text)
                    json_start = llm_response.find('{')
                    json_end = llm_response.rfind('}') + 1
                    if json_start != -1 and json_end != -1:
                        json_str = llm_response[json_start:json_end]
                        parsed_response = json.loads(json_str)
                        logger.info("Successfully parsed JSON response from LLM")
                        return parsed_response
                    else:
                        logger.warning("No JSON structure found in LLM response")
                except json.JSONDecodeError as e:
                    logger.warning(f"LLM response is not valid JSON: {e}")
                    # Try to find and extract JSON more carefully
                    try:
                        # Look for the first complete JSON object
                        brace_count = 0
                        start_idx = -1
                        for i, char in enumerate(llm_response):
                            if char == '{':
                                if start_idx == -1:
                                    start_idx = i
                                brace_count += 1
                            elif char == '}':
                                brace_count -= 1
                                if brace_count == 0 and start_idx != -1:
                                    json_str = llm_response[start_idx:i+1]
                                    parsed_response = json.loads(json_str)
                                    logger.info("Successfully parsed JSON using brace counting")
                                    return parsed_response
                    except json.JSONDecodeError:
                        logger.warning("Failed to parse JSON even with brace counting")
                
                # Fallback if JSON parsing fails - create structured content from raw response
                logger.warning("Using fallback content generation due to JSON parsing failure")
                
                # Try to extract meaningful content from the raw response
                if len(llm_response) > 1000:
                    fallback_notes = f"""Title: Document Analysis from {file_type}

I. Overview
- Content extracted from {file_type} document
- Key information and concepts identified

II. Key Concepts
- Important topics and definitions from the document

III. Summary
{llm_response[:1000]}..."""
                else:
                    fallback_notes = f"""Title: Document Analysis from {file_type}

I. Overview
- Content extracted from {file_type} document
- Key information and concepts identified

II. Key Concepts
- Important topics and definitions from the document

III. Summary
{llm_response}"""
                
                return {
                    "title": f"Document Analysis from {file_type}",
                    "summary": llm_response[:200] + "..." if len(llm_response) > 200 else llm_response,
                    "key_topics": ["Document Analysis", "Key Concepts"],
                    "definitions": [],
                    "study_notes": fallback_notes,
                    "difficulty": "Medium",
                    "subject": "General"
                }
            else:
                logger.error(f"LLM API error: {response.status_code}")
                return None
                
        except Exception as e:
            logger.error(f"Error processing with LLM: {e}")
            return None
    
    def save_to_database(self, file_path: str, content: str, file_type: str, llm_result: Dict, content_type: str = "both", user_id: str = None) -> bool:
        """Save processed content to database based on content type"""
        try:
            if not user_id:
                logger.error("User ID is required for saving content")
                return False
                
            if not self.db_connection or not self.db_connection.is_connected():
                if not self.connect_database():
                    return False
            
            cursor = self.db_connection.cursor()
            
            # Generate unique IDs
            note_id = str(uuid.uuid4())
            deck_id = str(uuid.uuid4())
            
            # Extract file name without extension
            file_name = Path(file_path).stem
            
            # Create note entry only if content_type includes notes
            if content_type in ["notes", "both"]:
                note_query = """
                INSERT INTO notes (id, user_id, title, content, category, created_at, modified_at, is_pinned)
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
                """
                note_values = (
                    note_id,
                    user_id,
                    llm_result.get('title', file_name),
                    llm_result.get('study_notes', content[:20000]),
                    llm_result.get('subject', 'General'),
                    datetime.now(),
                    datetime.now(),
                    False
                )
                cursor.execute(note_query, note_values)
                logger.info(f"Created note: {note_id}")
            
            # Create flashcard deck only if content_type includes flashcards
            if content_type in ["flashcards", "both"]:
                deck_query = """
                INSERT INTO flashcard_decks (id, user_id, title, description, subject, difficulty, created_at, last_studied, total_study_sessions)
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s)
                """
                deck_values = (
                    deck_id,
                    user_id,
                    f"{llm_result.get('title', file_name)} - Imported from {file_type}",
                    llm_result.get('summary', f"Content imported from {file_type} file"),
                    llm_result.get('subject', 'General'),
                    llm_result.get('difficulty', 'MEDIUM'),
                    datetime.now(),
                    None,
                    0
                )
                cursor.execute(deck_query, deck_values)
                logger.info(f"Created deck: {deck_id}")
            
            # Create flashcards from structured LLM response only if content_type includes flashcards
            if content_type in ["flashcards", "both"]:
                # First, try to use the new structured flashcards from LLM
                structured_flashcards = llm_result.get('flashcards', [])
                if structured_flashcards:
                    for flashcard_data in structured_flashcards:
                        flashcard_id = str(uuid.uuid4())
                        
                        flashcard_query = """
                        INSERT INTO flashcards (id, deck_id, question, answer, difficulty, created_at, last_studied, times_studied, is_correct)
                        VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s)
                        """
                        flashcard_values = (
                            flashcard_id,
                            deck_id,
                            flashcard_data.get('question', ''),
                            flashcard_data.get('answer', ''),
                            llm_result.get('difficulty', 'MEDIUM'),
                            datetime.now(),
                            None,
                            0,
                            False
                        )
                        cursor.execute(flashcard_query, flashcard_values)
                else:
                    # Fallback to old method if structured flashcards not available
                    definitions = llm_result.get('definitions', [])
                    for i, definition in enumerate(definitions[:25]):  # Increased limit for more flashcards
                        if ':' in definition:
                            term, definition_text = definition.split(':', 1)
                            flashcard_id = str(uuid.uuid4())
                            
                            flashcard_query = """
                            INSERT INTO flashcards (id, deck_id, question, answer, difficulty, created_at, last_studied, times_studied, is_correct)
                            VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s)
                            """
                            flashcard_values = (
                                flashcard_id,
                                deck_id,
                                f"What is {term.strip()}?",
                                definition_text.strip(),
                                llm_result.get('difficulty', 'MEDIUM'),
                                datetime.now(),
                                None,
                                0,
                                False
                            )
                            cursor.execute(flashcard_query, flashcard_values)
                
                    # Create flashcards from key topics as fallback
                    key_topics = llm_result.get('key_topics', [])
                    for i, topic in enumerate(key_topics[:15]):  # Increased limit for more topic flashcards
                        flashcard_id = str(uuid.uuid4())
                    
                        flashcard_query = """
                        INSERT INTO flashcards (id, deck_id, question, answer, difficulty, created_at, last_studied, times_studied, is_correct)
                        VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s)
                        """
                        flashcard_values = (
                            flashcard_id,
                            deck_id,
                            f"Explain: {topic}",
                            f"This is a key topic from the {file_type} document: {topic}",
                            llm_result.get('difficulty', 'MEDIUM'),
                            datetime.now(),
                            None,
                            0,
                            False
                        )
                        cursor.execute(flashcard_query, flashcard_values)
            
            self.db_connection.commit()
            cursor.close()
            
            logger.info(f"Successfully saved processed content to database")
            if content_type in ["notes", "both"]:
                logger.info(f"Created note: {note_id}")
            if content_type in ["flashcards", "both"]:
                logger.info(f"Created deck: {deck_id}")
                structured_flashcards = llm_result.get('flashcards', [])
                if structured_flashcards:
                    logger.info(f"Created {len(structured_flashcards)} structured flashcards")
                else:
                    definitions = llm_result.get('definitions', [])
                    key_topics = llm_result.get('key_topics', [])
                    logger.info(f"Created {len(definitions[:25]) + len(key_topics[:15])} fallback flashcards")
            
            return True
            
        except Error as e:
            logger.error(f"Database error: {e}")
            return False
    
    def process_document(self, file_path: str, content_type: str = "both", user_id: str = None) -> Dict:
        """Main method to process a document
        
        Args:
            file_path: Path to the document file
            content_type: Type of content to create ("notes", "flashcards", or "both")
            user_id: ID of the user creating the content
        """
        if not user_id:
            return {"success": False, "error": "User ID is required"}
            
        logger.info(f"Processing document: {file_path} for content type: {content_type} for user: {user_id}")
        
        # Extract content
        content, file_type = self.extract_content(file_path)
        if not content:
            return {"success": False, "error": "Failed to extract content"}
        
        # Process with LLM
        llm_result = self.process_with_llm(content, file_type)
        if not llm_result:
            return {"success": False, "error": "LLM processing failed"}
        
        # Save to database based on content type
        if not self.save_to_database(file_path, content, file_type, llm_result, content_type, user_id):
            return {"success": False, "error": "Database save failed"}
        
        # Prepare response based on content type
        response = {
            "success": True,
            "summary": llm_result.get('summary', ''),
            "subject": llm_result.get('subject', 'General'),
            "difficulty": llm_result.get('difficulty', 'Medium')
        }
        
        if content_type in ["notes", "both"]:
            response["note_title"] = llm_result.get('title', 'Unknown')
            response["note_id"] = llm_result.get('title', 'Unknown')
            # Return the formatted study notes content, not raw JSON
            response["note_content"] = llm_result.get('study_notes', 'No content generated')
        
        if content_type in ["flashcards", "both"]:
            response["deck_title"] = f"{llm_result.get('title', 'Unknown')} - Imported from {file_type}"
            structured_flashcards = llm_result.get('flashcards', [])
            if structured_flashcards:
                response["flashcards_created"] = len(structured_flashcards)
                response["flashcards"] = structured_flashcards  # Return formatted flashcards
            else:
                # Create flashcards from definitions and key topics
                flashcards = []
                definitions = llm_result.get('definitions', [])
                key_topics = llm_result.get('key_topics', [])
                
                # Add flashcards from definitions
                for definition in definitions[:25]:
                    if ':' in definition:
                        term, definition_text = definition.split(':', 1)
                        flashcards.append({
                            "question": f"What is {term.strip()}?",
                            "answer": definition_text.strip()
                        })
                
                # Add flashcards from key topics
                for topic in key_topics[:15]:
                    flashcards.append({
                        "question": f"Explain {topic}",
                        "answer": f"Key concept: {topic}"
                    })
                
                response["flashcards_created"] = len(flashcards)
                response["flashcards"] = flashcards
        
        return response

def main():
    """Main function for command line usage"""
    parser = argparse.ArgumentParser(description='Process documents for Study Space')
    parser.add_argument('file_path', help='Path to the document file')
    parser.add_argument('--config', default='../../config/ai-processor/config.json', help='Configuration file path')
    parser.add_argument('--output', help='Output directory for processed files')
    
    args = parser.parse_args()
    
    # Initialize processor
    processor = DocumentProcessor(args.config)
    
    # Process the document
    result = processor.process_document(args.file_path)
    
    if result['success']:
        print(f"✅ Successfully processed document!")
        print(f"📝 Note created: {result['note_id']}")
        print(f"🃏 Deck created: {result['deck_id']}")
        print(f"📚 Flashcards created: {result['flashcards_created']}")
        print(f"📖 Subject: {result['subject']}")
        print(f"📊 Difficulty: {result['difficulty']}")
        print(f"📄 Summary: {result['summary']}")
    else:
        print(f"❌ Error processing document: {result['error']}")
        sys.exit(1)

if __name__ == "__main__":
    main()